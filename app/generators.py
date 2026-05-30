"""Build professionally styled .docx and .pptx reports from the content model.

Two paths per format:
  * "standard" — a designed template built programmatically (cover page, accent
    colours, styled headings/tables, page numbers, section-aware slides).
  * "custom"   — the caller supplies their own .docx / .pptx template whose
    theme, fonts and layouts we reuse, so the output matches their branding.
"""
from __future__ import annotations

from io import BytesIO
from typing import Any, Optional

from docx import Document as DocxDocument
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt as PptPt

# ---- Palette (shared) -------------------------------------------------------
ACCENT_HEX = "5B21B6"      # indigo
DARK_HEX = "2E1065"        # deep indigo
LIGHT_ROW_HEX = "F2EEFB"   # pale tint for banded rows
# Word colours (docx.shared.RGBColor)
ACCENT = RGBColor(0x5B, 0x21, 0xB6)
DARK = RGBColor(0x2E, 0x10, 0x65)
GRAY = RGBColor(0x6B, 0x72, 0x80)
# PowerPoint colours (pptx.dml.color.RGBColor — a distinct type)
P_ACCENT = PptRGB(0x5B, 0x21, 0xB6)
P_DARK = PptRGB(0x2E, 0x10, 0x65)
P_BODY = PptRGB(0x1F, 0x29, 0x37)
P_WHITE = PptRGB(0xFF, 0xFF, 0xFF)

# Slide-packing budget for the standard content slides. Bullets are split so a
# slide holds at most _MAX_LINES_PER_SLIDE wrapped lines (estimated from text
# length) and at most _MAX_BULLETS items — whichever limit is hit first. This,
# together with shrink-to-fit on the text frame, keeps text inside the box.
_CHARS_PER_LINE = 90    # ~chars that fit on one line at 18pt across the body box
_MAX_LINES_PER_SLIDE = 13
_MAX_BULLETS = 7
_CONTENT_FONT_PT = 18
_TITLE_FONT_PT = 28


# =============================================================================
# WORD
# =============================================================================
def build_docx(
    title: str,
    blocks: list[dict[str, Any]],
    language_name: str,
    date_str: str,
    template_path: Optional[str] = None,
) -> bytes:
    if template_path:
        return _build_docx_template(title, blocks, language_name, date_str, template_path)
    return _build_docx_standard(title, blocks, language_name, date_str)


def _build_docx_standard(title, blocks, language_name, date_str) -> bytes:
    doc = DocxDocument()
    _style_docx_base(doc)
    _docx_cover(doc, title, language_name, date_str)
    _docx_page_numbers(doc)

    for block in blocks:
        btype = block["type"]
        if btype == "heading":
            level = min(max(int(block.get("level", 1)), 1), 4)
            doc.add_heading(block["text"], level=level)
        elif btype == "paragraph":
            doc.add_paragraph(block["text"])
        elif btype == "table":
            _docx_styled_table(doc, block["rows"])

    return _save_docx(doc)


def _build_docx_template(title, blocks, language_name, date_str, template_path) -> bytes:
    """Reuse the user's template styling; append the translated content."""
    doc = DocxDocument(template_path)

    def add(style_name: str, text: str, fallback: str = "Normal"):
        try:
            return doc.add_paragraph(text, style=style_name)
        except KeyError:
            return doc.add_paragraph(text, style=fallback)

    add("Title", title)
    sub = add("Subtitle", f"Translated Report — {language_name} · {date_str}", "Normal")
    if sub.runs:
        sub.runs[0].italic = True

    for block in blocks:
        btype = block["type"]
        if btype == "heading":
            level = min(max(int(block.get("level", 1)), 1), 4)
            try:
                doc.add_heading(block["text"], level=level)
            except KeyError:
                add("Normal", block["text"])
        elif btype == "paragraph":
            add("Normal", block["text"])
        elif btype == "table":
            _docx_styled_table(doc, block["rows"])

    return _save_docx(doc)


def _style_docx_base(doc) -> None:
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    pf = normal.paragraph_format
    pf.space_after = Pt(8)
    pf.line_spacing = 1.15

    for name, size, color in (
        ("Heading 1", 16, ACCENT),
        ("Heading 2", 13, DARK),
        ("Heading 3", 12, DARK),
    ):
        try:
            st = doc.styles[name]
        except KeyError:
            continue
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = color
        st.paragraph_format.space_before = Pt(14)
        st.paragraph_format.space_after = Pt(4)


def _docx_cover(doc, title, language_name, date_str) -> None:
    for _ in range(6):
        doc.add_paragraph()

    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run(title)
    run.bold = True
    run.font.size = Pt(30)
    run.font.color.rgb = DARK

    rule = doc.add_paragraph()
    rule.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _docx_bottom_border(rule, ACCENT_HEX, size=18)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Translated Report")
    r.italic = True
    r.font.size = Pt(14)
    r.font.color.rgb = GRAY

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    m = meta.add_run(f"{language_name}   ·   {date_str}")
    m.font.size = Pt(11)
    m.font.color.rgb = GRAY

    doc.add_page_break()


def _docx_bottom_border(paragraph, color_hex, size=6) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), str(size))
    bottom.set(qn("w:space"), "4")
    bottom.set(qn("w:color"), color_hex)
    pBdr.append(bottom)
    pPr.append(pBdr)


def _docx_shade_cell(cell, color_hex) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), color_hex)
    tcPr.append(shd)


def _docx_styled_table(doc, rows) -> None:
    if not rows:
        return
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=0, cols=cols)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"

    for ri, row in enumerate(rows):
        cells = table.add_row().cells
        for ci in range(cols):
            cell = cells[ci]
            cell.text = row[ci] if ci < len(row) else ""
            para = cell.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run("")
            run.font.size = Pt(10)
            if ri == 0:
                _docx_shade_cell(cell, ACCENT_HEX)
                run.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif ri % 2 == 0:
                _docx_shade_cell(cell, LIGHT_ROW_HEX)
    doc.add_paragraph()


def _docx_page_numbers(doc) -> None:
    footer = doc.sections[0].footer
    para = footer.paragraphs[0]
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = para.add_run()
    for tag, attr, val, text in (
        ("w:fldChar", "w:fldCharType", "begin", None),
        ("w:instrText", "xml:space", "preserve", "PAGE"),
        ("w:fldChar", "w:fldCharType", "end", None),
    ):
        el = OxmlElement(tag)
        el.set(qn(attr), val)
        if text:
            el.text = text
        run._r.append(el)


def _save_docx(doc) -> bytes:
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# =============================================================================
# POWERPOINT
# =============================================================================
def build_pptx(
    title: str,
    blocks: list[dict[str, Any]],
    language_name: str,
    date_str: str,
    template_path: Optional[str] = None,
) -> bytes:
    if template_path:
        return _build_pptx_template(title, blocks, language_name, date_str, template_path)
    return _build_pptx_standard(title, blocks, language_name, date_str)


def _estimate_lines(text: str) -> int:
    """Estimate how many wrapped lines a bullet occupies on a content slide."""
    return max(1, -(-len(text) // _CHARS_PER_LINE))  # ceil division


def _pack_bullets(bullets: list[str]) -> list[list[str]]:
    """Split bullets into slide-sized groups by estimated line count.

    A single very long bullet that exceeds the budget on its own still gets its
    own slide (shrink-to-fit on the text frame handles the remaining overflow).
    """
    groups: list[list[str]] = []
    current: list[str] = []
    used = 0
    for text in bullets:
        lines = _estimate_lines(text)
        too_tall = used + lines > _MAX_LINES_PER_SLIDE
        too_many = len(current) >= _MAX_BULLETS
        if current and (too_tall or too_many):
            groups.append(current)
            current, used = [], 0
        current.append(text)
        used += lines
    if current:
        groups.append(current)
    return groups


def _grouped_slides(title, blocks):
    """Yield ("content", slide_title, bullets) / ("table", slide_title, rows) /
    ("divider", heading, None) tuples, grouping paragraphs under their heading.
    """
    current_title = title
    bullets: list[str] = []
    rendered = True  # the title slide already represents the document title

    def flush():
        nonlocal bullets, rendered
        if not bullets:
            return
        first = True
        for chunk in _pack_bullets(bullets):
            slide_title = current_title if first else f"{current_title} (cont.)"
            yield_target.append(("content", slide_title, chunk))
            first = False
        bullets = []
        rendered = True

    yield_target: list = []
    for block in blocks:
        btype = block["type"]
        if btype == "heading":
            flush()
            if not rendered and current_title != title:
                yield_target.append(("divider", current_title, None))
            current_title = block["text"]
            rendered = False
        elif btype == "paragraph":
            bullets.append(block["text"])
        elif btype == "table":
            flush()
            yield_target.append(("table", current_title, block["rows"]))
            rendered = True
    flush()
    if not rendered and current_title != title:
        yield_target.append(("divider", current_title, None))
    return yield_target


def _build_pptx_standard(title, blocks, language_name, date_str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    # ---- Title slide ----
    slide = prs.slides.add_slide(blank)
    _fill_rect(slide, 0, 0, W, H, "FFFFFF")
    _fill_rect(slide, 0, 0, Inches(0.45), H, ACCENT_HEX)            # left bar
    _fill_rect(slide, 0, H - Inches(1.6), W, Inches(1.6), DARK_HEX)  # bottom band

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), W - Inches(2.0), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(40)
    p.font.bold = True
    p.font.color.rgb = P_DARK
    sp = tf.add_paragraph()
    sp.text = "Translated Report"
    sp.font.size = PptPt(22)
    sp.font.color.rgb = P_ACCENT

    meta = slide.shapes.add_textbox(Inches(1.0), H - Inches(1.25), W - Inches(2.0), Inches(0.9))
    mp = meta.text_frame.paragraphs[0]
    mp.text = f"{language_name}    ·    {date_str}"
    mp.font.size = PptPt(16)
    mp.font.color.rgb = P_WHITE

    # ---- Content / table / divider slides ----
    for kind, slide_title, payload in _grouped_slides(title, blocks):
        if kind == "content":
            _pptx_content_slide(prs, blank, W, H, slide_title, payload)
        elif kind == "table":
            _pptx_table_slide(prs, blank, W, H, slide_title, payload)
        elif kind == "divider":
            _pptx_divider_slide(prs, blank, W, H, slide_title)

    return _save_pptx(prs)


def _pptx_header(slide, W, title) -> None:
    _fill_rect(slide, 0, 0, W, Inches(0.18), ACCENT_HEX)  # top accent strip
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), W - Inches(1.2), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PptPt(_TITLE_FONT_PT)
    p.font.bold = True
    p.font.color.rgb = P_DARK


def _pptx_content_slide(prs, blank, W, H, title, bullets) -> None:
    slide = prs.slides.add_slide(blank)
    _pptx_header(slide, W, title)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), W - Inches(1.4), H - Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    # Shrink text to fit the box if a slide's content is still a little long.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + text
        p.font.size = PptPt(_CONTENT_FONT_PT)
        p.font.color.rgb = P_BODY
        p.space_after = PptPt(10)


def _pptx_table_slide(prs, blank, W, H, title, rows) -> None:
    slide = prs.slides.add_slide(blank)
    _pptx_header(slide, W, title)
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.7), Inches(1.7), W - Inches(1.4), Inches(0.4)
    ).table
    for r, row in enumerate(rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            para = cell.text_frame.paragraphs[0]
            para.font.size = PptPt(12)
            if r == 0:
                para.font.bold = True


def _pptx_divider_slide(prs, blank, W, H, title) -> None:
    slide = prs.slides.add_slide(blank)
    _fill_rect(slide, 0, 0, W, H, DARK_HEX)
    box = slide.shapes.add_textbox(Inches(1.0), H / 2 - Inches(0.8), W - Inches(2.0), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = PptPt(34)
    p.font.bold = True
    p.font.color.rgb = P_WHITE


def _fill_rect(slide, left, top, width, height, color_hex) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGB.from_string(color_hex)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _build_pptx_template(title, blocks, language_name, date_str, template_path) -> bytes:
    prs = Presentation(template_path)

    # Title slide via the template's own title layout.
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Translated Report — {language_name} · {date_str}"
    except (IndexError, AttributeError):
        pass

    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    for kind, slide_title, payload in _grouped_slides(title, blocks):
        if kind == "table":
            _tpl_table_slide(prs, content_layout, slide_title, payload)
        else:  # content or divider both render as a titled bullet slide
            bullets = payload if kind == "content" else []
            _tpl_content_slide(prs, content_layout, slide_title, bullets)

    return _save_pptx(prs)


def _tpl_content_slide(prs, layout, title, bullets) -> None:
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = title
    except AttributeError:
        pass
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:  # not the title placeholder
            body = ph
            break
    if body is None:
        return
    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text


def _tpl_table_slide(prs, layout, title, rows) -> None:
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = title
    except AttributeError:
        pass
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.7), Inches(1.7), Inches(8.5), Inches(0.4)
    ).table
    for r, row in enumerate(rows):
        for c in range(n_cols):
            table.cell(r, c).text = row[c] if c < len(row) else ""


def _save_pptx(prs) -> bytes:
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
